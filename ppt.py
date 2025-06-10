from pptx import Presentation
from pptx.util import Inches

# Create a new presentation
prs = Presentation()

# Define a light color theme (soft blue and white)
background_color = RGBColor(220, 240, 255)
title_color = RGBColor(0, 51, 102)

# Function to style slides with a lighter theme
def style_slide_light(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = background_color
    title = slide.shapes.title
    if title:
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = title_color
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Nanomaterials in Pharmaceuticals and Drug Development"
slide.placeholders[1].text = "Transforming Drug Delivery, Formulation, and Targeting"
style_slide_light(slide)

# Organized Topics
topics = [
    ("Introduction", "Nanomaterials revolutionize pharmaceuticals by enhancing drug delivery, solubility, and targeted treatment."),
    ("1) Drug Delivery", "Nanoparticles enable targeted drug delivery, improving efficacy and reducing side effects."),
    ("2) Enhanced Solubility", "Nanomaterials improve the solubility of poorly water-soluble drugs, increasing bioavailability."),
    ("3) Controlled Release", "Nanomaterial-based drug systems provide controlled and sustained drug release for better patient compliance."),
    ("4) Combination Therapies", "Nanomaterials facilitate multi-drug delivery, allowing effective combination therapies."),
    ("5) Personalized Medicine", "Nanomaterials help create patient-specific treatments for better outcomes."),
    ("6) Targeted Cancer Therapy", "Nanoparticles selectively accumulate in tumors, enhancing cancer treatment."),
    ("7) Immunotherapy", "Nanomaterials enhance the immune system’s ability to fight diseases."),
    ("8) Vaccine Delivery", "Nanoparticle-based vaccines improve stability, efficacy, and immune response."),
    ("9) Antimicrobial Agents", "Nanomaterials exhibit antimicrobial properties to fight drug-resistant infections."),
    ("10) Pharmacokinetics & Toxicology", "Nanomaterials assist in studying drug distribution, metabolism, and toxicity."),
    ("Conclusion", "Nanomaterials are revolutionizing medicine with effective and personalized drug treatments.")
]

# Create slides with the new theme
for title_text, content_text in topics:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = content_text
    style_slide_light(slide)

# Save the final version
final_pptx_path = "/mnt/data/Nanomaterials_in_Pharmaceuticals_Final.pptx"
prs.save(final_pptx_path)

final_pptx_path
